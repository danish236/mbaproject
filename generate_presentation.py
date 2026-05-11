import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Create presentation
prs = Presentation()

# Define Corporate Colors
asml_blue = RGBColor(0, 32, 91)
asml_orange = RGBColor(255, 102, 0)
dark_grey = RGBColor(64, 64, 64)

# ---------------------------------------------------------
# SLIDE 1: TITLE SLIDE
# ---------------------------------------------------------
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(3))
shape.fill.solid()
shape.fill.fore_color.rgb = asml_blue
shape.line.color.rgb = asml_blue

txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "A Dual-Horizon Financial Assessment"
p.font.bold = True
p.font.size = Pt(32)
p.font.name = 'Calibri'
p.font.color.rgb = RGBColor(255, 255, 255)

p2 = tf.add_paragraph()
p2.text = "Fundamental Resiliency and Long-Term Market Valuation of ASML Holding N.V."
p2.font.size = Pt(20)
p2.font.name = 'Calibri'
p2.font.color.rgb = asml_orange

txBox2 = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "Danish Shaikh | SAP ID: 77124115940\nMaster of Business Administration (Finance Specialization)\nMay 2026"
p3.font.size = Pt(18)
p3.font.name = 'Calibri'
p3.font.bold = True
p3.font.color.rgb = dark_grey


# ---------------------------------------------------------
# HELPER FUNCTION FOR DENSE CONTENT SLIDES
# ---------------------------------------------------------
def create_slide(title_text, points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = asml_blue
    title.text_frame.paragraphs[0].font.name = 'Calibri'
    title.text_frame.paragraphs[0].font.bold = True

    body = slide.placeholders[1]
    tf = body.text_frame
    for i, pt in enumerate(points):
        if i == 0:
            tf.text = pt
            tf.paragraphs[0].font.size = Pt(14)  # Reduced size for denser text
            tf.paragraphs[0].font.name = 'Calibri'
        else:
            p = tf.add_paragraph()
            p.text = pt
            p.font.size = Pt(14)  # Reduced size for denser text
            p.font.name = 'Calibri'
            p.space_before = Pt(8)
    return slide


# ---------------------------------------------------------
# CONTENT SLIDES (PULLED FROM THESIS)
# ---------------------------------------------------------
create_slide("1. Executive Summary & Objective", [
    "Objective: To empirically deconstruct the financial resiliency and market valuation of ASML[cite: 374, 376].",
    "Methodology: An automated, Python-driven extraction pipeline analyzing a 4-year fundamental window and a 15-year technical horizon[cite: 377, 378].",
    "Core Hypothesis: ASML's absolute hardware monopoly insulates it from standard corporate finance vulnerabilities[cite: 403, 1158].",
    "Outcome: A stark dichotomy was found between its flawless internal fundamentals and its severely overvalued public equity pricing[cite: 385, 386]."
])

create_slide("2. The EUV Monopoly & Economic Moat", [
    "ASML operates as the exclusive global supplier of Extreme Ultraviolet (EUV) and High-NA EUV lithography systems[cite: 397].",
    "Commanding an 82.5% market share, it effectively functions as a de facto global monopoly, defeating legacy competitors Nikon (10.2%) and Canon (6.1%)[cite: 375, 563, 564].",
    "Inelastic Pricing Power: The inability of foundries (TSMC, Intel, Samsung) to source alternative machinery allows ASML to pass inflationary pressures to clients[cite: 402, 652].",
    "Structural Moat: Protected by thousands of patents, specialized quantum engineering talent, and exclusive supplier alliances (e.g., Carl Zeiss)[cite: 462, 502].",
    "The 18-month manufacturing lead time creates a multi-year order backlog exceeding $30 billion, shielding ASML from standard supply chain 'Bullwhip' effects[cite: 422, 423]."
])

create_slide("3. Fundamental Health: Margin Expansion", [
    "Total Revenue surged from $21.5 billion in 2021 to $29.8 billion in 2023, representing an absolute top-line expansion of nearly 40% over 24 months[cite: 614, 615].",
    "DuPont Analysis confirms extraordinary Net Profit Margins: expanding from 26.56% in 2022 to a peak of 29.41% in 2025[cite: 651].",
    "Gross Margins consistently clear the 52% threshold, proving the physical cost of a $350 million machine consumes less than half its sale price[cite: 677, 1111].",
    "R&D Intensity: ASML routinely allocates upwards of 14.5% to 17% of revenue ($4.0 billion in 2023) to R&D, continuously expanding its barrier to entry[cite: 642, 643, 713].",
    "Statutory Tax Optimization via the Dutch 'Innovation Box' regime lowers the effective tax rate to ~15%, drastically inflating Net Operating Profit After Tax[cite: 807, 808, 809, 812]."
])

create_slide("4. Solvency & Working Capital Masterclass", [
    "Deleveraging: Debt-to-Equity ratio systematically compressed from 0.48 in 2022 to a deeply conservative 0.22 in 2025, nullifying interest rate risk[cite: 680].",
    "Altman Z-Score oscillates between 12.0 and 15.0, proving insolvency distress is statistically indistinguishable from zero[cite: 820, 823].",
    "Standard Current Ratios are misleading due to massive Work-in-Progress (WIP) inventory, but the Quick (Acid-Test) Ratio confirms pristine immediate liquidity[cite: 682, 683, 688, 690].",
    "ASML mitigates WIP drag via a negative working capital structure, forcing foundries to front-load massive, non-refundable down payments[cite: 768, 770, 772].",
    "This supplier-financed architecture frees up internal cash reserves for aggressive share buybacks and dividend compounding[cite: 772, 773, 901]."
])

create_slide("5. Long-Term Market Valuation & Momentum", [
    "The 15-Year Compound Annual Growth Rate (CAGR) stands at an anomalous 28.00%, vastly outperforming standard macroeconomic benchmarks[cite: 845, 846].",
    "Equity ownership is extremely concentrated: exactly 80.0% of the float is tightly held by global institutional asset managers, accelerating price momentum[cite: 873].",
    "A calculated CAPM Beta of 1.82 mathematically proves the asset is 82% more volatile than the broader S&P 500, indicating high systematic risk[cite: 880].",
    "Maximum Drawdown profiling reveals a -50% peak-to-trough contraction during the 2022 central bank tightening cycle, highlighting extreme interest-rate sensitivity[cite: 865, 866].",
    "Free Cash Flow to Equity (FCFE) and Free Cash Flow to Firm (FCFF) mathematically converge due to negligible debt servicing costs[cite: 939, 940]."
])

# ---------------------------------------------------------
# SLIDE 6: THE VALUATION DISCONNECT (AUTO-CHART & TEXT)
# ---------------------------------------------------------
slide6 = prs.slides.add_slide(prs.slide_layouts[5])
title6 = slide6.shapes.title
title6.text = "6. Intrinsic Valuation: The DCF Disconnect"
title6.text_frame.paragraphs[0].font.color.rgb = asml_blue
title6.text_frame.paragraphs[0].font.bold = True
title6.text_frame.paragraphs[0].font.name = 'Calibri'

chart_data = CategoryChartData()
chart_data.categories = ['Intrinsic Value (DCF Model)', 'Current Market Price']
chart_data.add_series('USD', (411.34, 1592.02))

x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5)
chart = slide6.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
chart.has_legend = False

# Add dense text next to chart
tx = slide6.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(4.5))
tf = tx.text_frame
tf.word_wrap = True
tf.text = "A 5-Year DCF model was executed using an 8.5% WACC, 15% near-term growth, and 2.5% terminal growth[cite: 949]."
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.name = 'Calibri'

p = tf.add_paragraph()
p.text = "The model generated an intrinsic share value of exactly $411.34, starkly contrasting with current trading prices exceeding $1,592.02[cite: 960, 961]."
p.font.size = Pt(14);
p.font.name = 'Calibri';
p.space_before = Pt(8)

p2 = tf.add_paragraph()
p2.text = "The public equities market assigns a massive 'Scarcity Premium,' pricing in a flawless, multi-decade AI super-cycle with zero margin of safety[cite: 962, 983, 984]."
p2.font.size = Pt(14);
p2.font.name = 'Calibri';
p2.space_before = Pt(8)

p3 = tf.add_paragraph()
p3.text = "Even highly optimistic 'Bull Market' Monte Carlo parameters (7.5% WACC, 3.5% Terminal Growth) max out at a $580.00 intrinsic value[cite: 978, 979, 982]."
p3.font.size = Pt(14);
p3.font.name = 'Calibri';
p3.space_before = Pt(8)

# ---------------------------------------------------------
# CONTINUED CONTENT SLIDES
# ---------------------------------------------------------
create_slide("7. Relative Valuation & Real Options", [
    "Comparable Company Analysis (CCA) shows ASML trading at a Forward P/E of 40x-45x, double the 18x-22x baseline of peers like Lam Research[cite: 1007, 1008, 1009].",
    "This extreme multiple expansion leaves the asset highly vulnerable to 'multiple compression' if foundry CapEx decelerates or operational missteps occur[cite: 1014, 1015].",
    "Real Options Valuation (ROV) suggests the market is not just valuing current deliveries, but pricing the High-NA R&D pipeline as a massive financial call option[cite: 1020, 1021, 1022, 1030].",
    "Economic Value Added (EVA) analysis proves a 30% positive spread between ASML's 40%+ ROIC and its 8.5% WACC[cite: 1036, 1037, 1038]."
])

create_slide("8. The Geopolitical Black Swan", [
    "The ongoing US-China Semiconductor Cold War presents a severe, unquantifiable threat to ASML’s Total Addressable Market (TAM)[cite: 1049, 1050, 1179].",
    "US Foreign Direct Product Rule (FDPR) extraterritorial controls bar ASML from shipping flagship EUV and legacy DUV systems to Chinese entities[cite: 1054, 1055].",
    "This macroeconomic geographic concentration risk threatens the 15-20% of global system sales historically derived from the Chinese mainland[cite: 1060, 1061, 1062].",
    "Macro-regulatory tailwinds, such as the $52B U.S. CHIPS Act and €43B European Chips Act, partially underwrite client CapEx and offset export bans[cite: 1066, 1067, 1070]."
])

create_slide("9. The High-NA EUV Financial Runway", [
    "Future revenue velocity is anchored by the Twinscan EXE (High-NA EUV) series, raising the Average Selling Price (ASP) to $350 million[cite: 1108, 1111].",
    "This pricing architecture allows ASML to achieve 15% top-line growth without requiring proportional increases in physical manufacturing throughput[cite: 1113, 1116].",
    "Installed Base Management (Services) acts as a high-margin, 'Razor and Blade' software annuity, stabilizing cash flow during CapEx downturns[cite: 1119, 1122, 1123].",
    "Corporate treasury utilizes highly sophisticated FX derivative overlays to insulate this global revenue stream from translational and transactional currency shocks[cite: 1136, 1137, 1138]."
])

create_slide("10. Final Strategic Recommendation", [
    "FOR EXISTING SHAREHOLDERS: HOLD. The pristine corporate finance architecture and multi-decade AI CapEx cycle justify retaining exposure[cite: 1173].",
    "The structural alignment of the Dutch Two-Tier Board and ROAIC-based executive compensation mitigates agency risk and ensures capital discipline[cite: 1146, 1147, 1153].",
    "FOR NEW CAPITAL DEPLOYMENT: AVOID. The extreme valuation premium ($1,592 vs $411 DCF) offers an unacceptable level of systematic risk[cite: 1174].",
    "Allocators must await a localized macroeconomic recession or geopolitical shock to trigger a 30% to 40% equity drawdown for a mathematically defensible entry point[cite: 1086]."
])

prs.save('ASML_Dense_Master_Deck.pptx')
print("Dense 10-Slide Deck Generated!")